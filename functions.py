import urllib.request
import requests
import io
from PIL import Image
from bs4 import BeautifulSoup
from pptx import Presentation  

headers_Get = {'User-Agent': 'Mozilla/5.0 (Windows NT 6.1; WOW64; rv:49.0) Gecko/20100101 Firefox/49.0', 'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
'Accept-Language': 'en-US,en;q=0.5','Accept-Encoding': 'gzip, deflate','DNT': '1','Connection': 'keep-alive','Upgrade-Insecure-Requests': '1'}

def yahoo_pictures(q):
    s = requests.Session()
    q = '+'.join(q.split())
    url = f'https://images.search.yahoo.com/search/images;_ylt=AwrE1xifo.xfcrUAxwRXNyoA;_ylu=Y29sbwNiZjEEcG9zAzEEdnRpZAMEc2VjA3BpdnM-?p={q}&fr2=piv-web&fr=yfp-t'
    r = s.get(url, headers=headers_Get)
    soup = BeautifulSoup(r.text, "html.parser")
    output = []
    label_log = []
    url_log = []
    for searchWrapper in soup.find_all('ul', {'id':'sres'}):
        a = searchWrapper.find_all('a')
        img = searchWrapper.find_all('img')
        for i, this_img in enumerate(img[::2]):
            img_src = this_img["data-src"]
            url_log.append(img_src)
        for el in a:
            label_log.append(el["aria-label"])
        for i in range(min(len(label_log), len(url_log))):
            result = {'url' : url_log[i], 'label' : label_log[i]}
            output.append(result)
    return output 
    
def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]

    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size

    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width

    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)

    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side
        
def add_picture_slide(prs, url, image_title=' ', image_subtitle=' '):
    page = requests.get(url)

    image_data = page.content # byte values of the image
    image = Image.open(io.BytesIO(image_data))

    layout8 = prs.slide_layouts[8]
    slide = prs.slides.add_slide(layout8)

    title = slide.shapes.title.text = image_title
    sub = slide.placeholders[2].text = image_subtitle
    _add_image(slide,1, io.BytesIO(image_data))

def ask_exit(s=''):
    print()
    print(s)
    print('Вы хотите продолжить? (д/н)')
    return str(input("    --> ")) != 'д'

def search_web(topic, pictures_count):
    return yahoo_pictures(topic)[:pictures_count]

def multi_split(s):
    arr = []
    for row in s.split('|'):
        for el in row.split(':'):
            arr = arr + el.split('-')
    if len(arr) == 2:
        arr =  arr[1:]
    elif len(arr) == 3:
        arr =  arr[1:-1]
    sum = ''
    for el in arr:
        sum = sum + el
    return sum
    
def ask_session():
    topic = str(input("""Введите тему презентации
    --> """))
    pictures_count = int(input("""Введите число слайдов
    --> """))-1
    ppt = Presentation() 
    first_slide_layout = ppt.slide_layouts[0]  

    slide = ppt.slides.add_slide(first_slide_layout) 
    slide.shapes.title.text = topic
    search_results = search_web(topic, pictures_count)
    urls = [el['url'] for el in search_results]
    exit_flag = False
    if len(urls) < pictures_count:
        exit_flag = ask_exit(f"Вместо {pictures_count+1} слайдов будет только {len(urls)+1}")
    if exit_flag == True:
        print('Работа прервана ...')
        
        
    else:
#         for el in urls:
#             urls.append(el['url'])
        for i, url in enumerate(urls):
            add_picture_slide(ppt, url, image_title=' ', image_subtitle=multi_split(search_results[i]["label"]))
        ppt.save(topic + ".pptx") 
        print(f"\nСоздан файл '{topic + '.pptx'}'")