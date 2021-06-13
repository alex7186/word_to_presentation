from io import BytesIO
from PIL import Image
from bs4 import BeautifulSoup

from pptx import Presentation  
from requests import Session
from requests import get as req_get
from re import split as re_split

# строка для запроса
yahoo_url = 'https://images.search.yahoo.com/search/images;_ylt=AwrJ7JzuYMZguSYAetpXNyoA;_ylu=Y29sbwNiZjEEcG9zAzEEdnRpZANDMTg2NV8xBHNlYwNwaXZz?p={q}&fr2=piv-web&fr=yfp-t-s'

# заголовки 
headers_Get = {'User-Agent': 'Mozilla/5.0 (Windows NT 6.1; WOW64; rv:49.0) Gecko/20100101 Firefox/49.0', 'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
'Accept-Language': 'en-US,en;q=0.5','Accept-Encoding': 'gzip, deflate','DNT': '1','Connection': 'keep-alive','Upgrade-Insecure-Requests': '1'}

progress_bar_count=1700

def yahoo_pictures(question):
    req_session = Session()
    question = '+'.join(question.split())
    
    # получен ответ на запрос, парсинг string с помощью bs4
    current_request = req_session.get(yahoo_url.format(q=question), headers=headers_Get)
    soup = BeautifulSoup(current_request.text, "html.parser")

    # счетчики для url и описаний изображений
    label_log = []
    url_log = []

    output = []

    
    # i=0
    for i, searchWrapper in enumerate(soup.find_all('ul', {'id':'sres'})):

        a_elements = searchWrapper.find_all('a')
        img_elements = searchWrapper.find_all('img')


        for i, this_img in enumerate(img_elements[::2]):
            img_src = this_img["data-src"]
            url_log.append(img_src)

        # иногда в этом месте происходит AttributeError
        for el in a_elements:
            label_log.append(el["aria-label"])

        for ii in range(min(len(label_log), len(url_log))):
            result = {'url' : url_log[ii], 'label' : label_log[ii]}
            output.append(result)

    return output 


def add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]

    im = Image.open(image_url)
    width, height = im.size

    placeholder.height = height
    placeholder.width = width

    # добавляем изображение
    placeholder = placeholder.insert_picture(image_url)

    # вычисляем и сравниваем диагональные коэффициенты
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    # placeholder слишком широкий
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side

    # placeholder слишком высокий
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side
        

def add_picture_slide(prs, url, image_title=' ', image_subtitle=' '):
    page = req_get(url)

    image_data = page.content 
    image = Image.open(BytesIO(image_data)) # открываем изображение в байтах

    # добавляем и оформляем слайд
    layout8 = prs.slide_layouts[8] 
    slide = prs.slides.add_slide(layout8) 

    title = slide.shapes.title.text = image_title
    sub = slide.placeholders[2].text = image_subtitle
    add_image(slide,1, BytesIO(image_data))


def ask_exit(s=''):
    print()
    print(s)
    print('Вы хотите продолжить? (д/н)')
    return str(input("\t--> ")) != 'д'


def search_web(topic, pictures_count):
    return yahoo_pictures(topic)[:pictures_count]


def multi_split(s):
    # убираем разделители описания изображения из yahoo
    arr = re_split(', |_|-|!|:|;', s)
    arr = [el[1:] if el[0]==' ' else el for el in arr]
    return max(arr)


# https://stackoverflow.com/questions/3173320/text-progress-bar-in-the-console
# Print iterations progress 
def printProgressBar (iteration, total, prefix = '', suffix = '', decimals = 1, length = 100, fill = '█', printEnd = "\r"):
    """
    Call in a loop to create terminal progress bar
    @params:
        iteration   - Required  : current iteration (Int)
        total       - Required  : total iterations (Int)
        prefix      - Optional  : prefix string (Str)
        suffix      - Optional  : suffix string (Str)
        decimals    - Optional  : positive number of decimals in percent complete (Int)
        length      - Optional  : character length of bar (Int)
        fill        - Optional  : bar fill character (Str)
        printEnd    - Optional  : end character (e.g. "\r", "\r\n") (Str)
    """
    percent = ("{0:." + str(decimals) + "f}").format(100 * (iteration / float(total)))
    filledLength = int(length * iteration // total)
    bar = fill * filledLength + '-' * (length - filledLength)
    print(f'\r{prefix} |{bar}| {percent}% {suffix}', end = printEnd)
    # Print New Line on Complete
    if iteration == total: 
        print()

    
def ask_session(save_file=True, topic='', slides_count=0):
    # по информации от пользователя формируется pptx 
    if not topic:
        topic = str(input("""Введите тему презентации\n\t--> """))

    if not slides_count:
        slides_count = int(input("""Введите число слайдов\n\t--> """))-1

    ppt = Presentation() 
    first_slide_layout = ppt.slide_layouts[0]  

    slide = ppt.slides.add_slide(first_slide_layout) 
    slide.shapes.title.text = topic

    # при возникновении ошибки возвращаем False в флаге
    # (на стороне yahoo, не всегда изображения имеют требуемый класс) 
    printProgressBar(0, progress_bar_count, prefix = 'Progress:', suffix = 'Complete', length = 50)

    try:
        search_results = search_web(topic, slides_count)
    except:
        return False, topic, slides_count

    printProgressBar(int(progress_bar_count*0.39), progress_bar_count, prefix = 'Progress:', suffix = 'Complete', length = 50)
    urls = [el['url'] for el in search_results]

    # если изображений меньше, чем необходимо, задается вопрос о продолжении работы
    exit_flag = False
    if len(urls) < slides_count:
        exit_flag = ask_exit(f"Вместо {slides_count+1} слайдов будет только {len(urls)+1}")

    if exit_flag == True:
        printProgressBar(int(progress_bar_count*1), progress_bar_count, prefix = 'Progress:', suffix = 'Complete', length = 50)
        print('Работа прервана ...')
        
        
    # если изображений достаточно и необходимо сохранить файл:
    else:
        for i, url in enumerate(urls):
            add_picture_slide(ppt, url, image_title=' ', image_subtitle=search_results[i]["label"]) 
        if save_file:
            ppt.save(topic + ".pptx") 
        printProgressBar(int(progress_bar_count*1), progress_bar_count, prefix = 'Progress:', suffix = 'Complete', length = 50)
        print(f"Создан файл '{topic + '.pptx'}'")

    return True, topic, slides_count


res = ask_session()
if not res[0]:
    ask_session(topic=res[1], slides_count=res[2])
